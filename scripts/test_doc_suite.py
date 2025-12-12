#!/usr/bin/env python3
"""Generate long doc/docx/pptx files (>=50 pages, >=500 Chinese chars/page) and submit conversions.

Scenarios (per 10 items each):
- docx (50 pages) -> pdf
- doc (generated via soffice from docx) -> pdf
- pptx (50 slides) -> pdf (if service supports pptx->pdf)

Env:
- API_URL (default http://127.0.0.1:8000/api/v1/convert)
- API_APPID / API_KEY
- CONCURRENCY (default 5)
- PAGES (default 50)
- CHARS_PER_PAGE (default 500)

Note: doc generation uses local soffice to convert a docx into .doc. Requires LibreOffice CLI available.
"""

from __future__ import annotations

import base64
import json
import os
import tempfile
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from typing import Any, Dict, List
from urllib.error import HTTPError, URLError
from urllib.request import Request, urlopen

from docx import Document
from pptx import Presentation

API_URL = os.getenv("API_URL", "http://127.0.0.1:8000/api/v1/convert")
APP_ID = os.getenv("API_APPID", "12872b4f6d05")
API_KEY = os.getenv("API_KEY", "CDqH3KrwthT8GtX9TEQppLMOrB96N178zR-MzGgdEEk")
CONCURRENCY = int(os.getenv("CONCURRENCY", "5"))
PAGES = int(os.getenv("PAGES", "50"))
CHARS_PER_PAGE = int(os.getenv("CHARS_PER_PAGE", "500"))


CH_TEXT = "汉" * CHARS_PER_PAGE


def make_docx(path: Path, label: str) -> None:
    doc = Document()
    for i in range(PAGES):
        p = doc.add_paragraph(f"{label} 第{i+1}页 ")
        p.add_run(CH_TEXT)
        doc.add_page_break()
    doc.save(path)


def make_pptx(path: Path, label: str) -> None:
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for i in range(PAGES):
        slide = prs.slides.add_slide(blank)
        tx_box = slide.shapes.add_textbox(50, 50, 600, 400)
        tf = tx_box.text_frame
        tf.text = f"{label} 第{i+1}页 " + CH_TEXT
    prs.save(path)


def docx_to_doc(docx_path: Path, doc_path: Path) -> None:
    import subprocess

    cmd = [
        "soffice",
        "--headless",
        "--convert-to",
        "doc",
        "--outdir",
        str(doc_path.parent),
        str(docx_path),
    ]
    subprocess.run(cmd, check=True)
    generated = doc_path.parent / f"{docx_path.stem}.doc"
    generated.rename(doc_path)


def to_b64(path: Path) -> tuple[str, float]:
    data = path.read_bytes()
    b64 = base64.b64encode(data).decode("ascii")
    size_mb = round(len(data) / (1024 * 1024), 4)
    return b64, size_mb


def submit(payload: Dict[str, Any], name: str) -> Dict[str, Any]:
    data = json.dumps(payload).encode("utf-8")
    headers = {"Content-Type": "application/json", "X-Appid": APP_ID, "X-Key": API_KEY}
    req = Request(API_URL, data=data, headers=headers, method="POST")
    try:
        with urlopen(req, timeout=60) as resp:
            body = resp.read().decode("utf-8", errors="replace")
            return {"name": name, "status": "ok", "http": resp.status, "body": body}
    except HTTPError as exc:
        try:
            detail = exc.read().decode("utf-8", errors="replace")
        except Exception:
            detail = exc.reason
        return {"name": name, "status": "fail", "http": exc.code, "body": detail}
    except URLError as exc:
        return {"name": name, "status": "fail", "http": None, "body": str(exc)}


def build_payload(
    b64: str,
    size_mb: float,
    source_fmt: str,
    target_fmt: str,
    filename: str,
    *,
    page_limit: int | None = None,
) -> Dict[str, Any]:
    return {
        "task_name": f"{source_fmt}-to-{target_fmt}",
        "priority": "normal",
        "files": [
            {
                "base64_data": b64,
                "source_format": source_fmt,
                "target_format": target_fmt,
                "filename": filename,
                "size_mb": size_mb,
                "page_limit": page_limit,
            }
        ],
    }


def main() -> int:
    results: List[Dict[str, Any]] = []
    with tempfile.TemporaryDirectory() as tmpdir:
        tmp = Path(tmpdir)
        docx_files: List[Path] = []
        doc_files: List[Path] = []
        pptx_files: List[Path] = []

        for i in range(10):
            docx_path = tmp / f"docx_{i}.docx"
            make_docx(docx_path, "DOCX混合测试")
            docx_files.append(docx_path)

            doc_path = tmp / f"doc_{i}.doc"
            docx_to_doc(docx_path, doc_path)
            doc_files.append(doc_path)

            pptx_path = tmp / f"pptx_{i}.pptx"
            make_pptx(pptx_path, "PPTX混合测试")
            pptx_files.append(pptx_path)

        requests: List[tuple[str, Dict[str, Any]]] = []
        # full documents
        for path in doc_files:
            b64, size_mb = to_b64(path)
            requests.append(("doc->pdf", build_payload(b64, size_mb, "doc", "pdf", path.name)))
        for path in docx_files:
            b64, size_mb = to_b64(path)
            requests.append(("docx->pdf", build_payload(b64, size_mb, "docx", "pdf", path.name)))
        for path in pptx_files:
            b64, size_mb = to_b64(path)
            requests.append(("pptx->pdf", build_payload(b64, size_mb, "pptx", "pdf", path.name)))

        # page-limited documents (first few items only to keep load reasonable)
        PAGE_LIMIT = 5
        for path in doc_files[:3]:
            b64, size_mb = to_b64(path)
            requests.append(("doc->pdf (pages)", build_payload(b64, size_mb, "doc", "pdf", path.name, page_limit=PAGE_LIMIT)))
        for path in docx_files[:3]:
            b64, size_mb = to_b64(path)
            requests.append(("docx->pdf (pages)", build_payload(b64, size_mb, "docx", "pdf", path.name, page_limit=PAGE_LIMIT)))
        for path in pptx_files[:3]:
            b64, size_mb = to_b64(path)
            requests.append(("pptx->pdf (pages)", build_payload(b64, size_mb, "pptx", "pdf", path.name, page_limit=PAGE_LIMIT)))

        with ThreadPoolExecutor(max_workers=CONCURRENCY) as pool:
            future_map = {pool.submit(submit, payload, name): name for name, payload in requests}
            for fut in as_completed(future_map):
                res = fut.result()
                results.append(res)
                status = "OK" if res["status"] == "ok" else "FAIL"
                print(f"[{status}] {res['name']} -> {res['http']} {res['body'][:120]}")

    ok = sum(1 for r in results if r["status"] == "ok")
    total = len(results)
    print(f"Summary: {ok}/{total} succeeded")
    return 0 if ok == total else 1


if __name__ == "__main__":
    raise SystemExit(main())
