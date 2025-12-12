#!/usr/bin/env python3
"""Concurrent conversion load test covering image/audio/video mixes.

Scenarios (10 requests each by default):
- html(base64, ~50 pages) -> pdf
- svg(base64, padded) -> png
- wav(base64 silence, ~50MB) -> mp3
- gif(base64 padded, ~100MB) -> mp4

Usage:
    API_URL=http://127.0.0.1:8000/api/v1/convert \
    API_APPID=xxx API_KEY=yyy \
    python scripts/test_concurrent_conversions.py

Env vars:
- API_URL: conversion endpoint
- API_APPID / API_KEY: auth headers
- CONCURRENCY: max in-flight requests (default 10)
- IMAGE_SIZE_MB: per-image payload size (default 5)
- AUDIO_SIZE_MB: per-audio payload size (default 50)
- VIDEO_SIZE_MB: per-video payload size (default 100)
- DOC_PAGES: pages to generate in HTML (default 50)
"""

from __future__ import annotations

import base64
import json
import os
import time
from concurrent.futures import ThreadPoolExecutor, as_completed
from functools import lru_cache
from pathlib import Path
from tempfile import TemporaryDirectory
from typing import Any, Dict, List
from urllib.error import HTTPError, URLError
from urllib.request import Request, urlopen

from docx import Document
from pptx import Presentation

API_URL = os.getenv("API_URL", "http://127.0.0.1:8000/api/v1/convert")
APP_ID = os.getenv("API_APPID", "12872b4f6d05")
API_KEY = os.getenv("API_KEY", "CDqH3KrwthT8GtX9TEQppLMOrB96N178zR-MzGgdEEk")
CONCURRENCY = int(os.getenv("CONCURRENCY", "10"))
IMAGE_SIZE_MB = float(os.getenv("IMAGE_SIZE_MB", "5"))
AUDIO_SIZE_MB = float(os.getenv("AUDIO_SIZE_MB", "50"))
VIDEO_SIZE_MB = float(os.getenv("VIDEO_SIZE_MB", "100"))
DOC_PAGES = int(os.getenv("DOC_PAGES", "50"))
DOC_CHARS_PER_PAGE = int(os.getenv("DOC_CHARS_PER_PAGE", "500"))
REQUESTS_PER_SCENARIO = int(os.getenv("REQUESTS_PER_SCENARIO", "10"))
HTTP_TIMEOUT = float(os.getenv("HTTP_TIMEOUT", "30"))
RETRIES = int(os.getenv("RETRIES", "0"))
PAGE_LIMIT = int(os.getenv("PAGE_LIMIT", "5"))

SCENARIO_META = {
    "html(base64 50p)->pdf": "html->pdf",
    "html(base64 50p)->pdf (pages)": "html->pdf",
    "svg(base64)->png": "svg->png",
    "wav(base64)->mp3": "wav->mp3",
    "gif(base64)->mp4": "gif->mp4",
    "doc->pdf": "doc->pdf",
    "doc->pdf (pages)": "doc->pdf",
    "docx->pdf": "docx->pdf",
    "docx->pdf (pages)": "docx->pdf",
    "ppt->pdf": "ppt->pdf",
    "ppt->pdf (pages)": "ppt->pdf",
    "pptx->pdf": "pptx->pdf",
    "pptx->pdf (pages)": "pptx->pdf",
}


def _post(payload: Dict[str, Any], name: str) -> dict[str, Any]:
    data = json.dumps(payload).encode("utf-8")
    headers = {
        "Content-Type": "application/json",
        "X-Appid": APP_ID,
        "X-Key": API_KEY,
    }

    attempts = max(1, RETRIES + 1)
    last_error: dict[str, Any] | None = None

    for attempt in range(1, attempts + 1):
        req = Request(API_URL, data=data, headers=headers, method="POST")
        start = time.perf_counter()
        try:
            with urlopen(req, timeout=HTTP_TIMEOUT) as resp:
                body = resp.read().decode("utf-8", errors="replace")
                elapsed = time.perf_counter() - start
                return {"scenario": name, "status": "ok", "http": resp.status, "body": body, "elapsed": elapsed, "attempt": attempt}
        except HTTPError as exc:
            try:
                detail = exc.read().decode("utf-8", errors="replace")
            except Exception:
                detail = exc.reason
            elapsed = time.perf_counter() - start
            last_error = {"scenario": name, "status": "fail", "http": exc.code, "body": detail, "elapsed": elapsed, "attempt": attempt}
        except URLError as exc:
            elapsed = time.perf_counter() - start
            last_error = {"scenario": name, "status": "fail", "http": None, "body": str(exc), "elapsed": elapsed, "attempt": attempt}

    return last_error or {"scenario": name, "status": "fail", "http": None, "body": "unknown error", "elapsed": 0.0, "attempt": attempts}


@lru_cache(maxsize=None)
def _cached_svg_payload() -> tuple[str, float]:
    svg = f"""<svg xmlns='http://www.w3.org/2000/svg' width='640' height='640'><rect width='640' height='640' fill='teal'/><text x='32' y='360' font-size='160' fill='white'>KT</text></svg>"""
    raw = bytearray(svg.encode("utf-8"))
    target_bytes = int(IMAGE_SIZE_MB * 1024 * 1024)
    if len(raw) < target_bytes:
        raw.extend(b" <!--pad-->" * ((target_bytes - len(raw)) // 9 + 1))
        raw = raw[:target_bytes]
    b64 = base64.b64encode(raw).decode("ascii")
    size_mb = round(len(raw) / (1024 * 1024), 4)
    return b64, size_mb


def build_svg_payload(idx: int) -> Dict[str, Any]:
    b64, size_mb = _cached_svg_payload()
    return {
        "task_name": f"svg-to-png-{idx}",
        "priority": "normal",
        "files": [
            {
                "base64_data": b64,
                "source_format": "svg",
                "target_format": "png",
                "filename": f"logo-{idx}.svg",
                "size_mb": size_mb,
            }
        ],
    }


def build_html_payload(idx: int) -> Dict[str, Any]:
    b64, size_mb = _cached_html_payload()
    return {
        "task_name": f"html-to-pdf-{idx}",
        "priority": "normal",
        "files": [
            {
                "base64_data": b64,
                "source_format": "html",
                "target_format": "pdf",
                "filename": f"long-{idx}.html",
                "size_mb": size_mb,
            }
        ],
    }


def build_html_payload_pagelimit(idx: int) -> Dict[str, Any]:
    payload = build_html_payload(idx)
    payload["files"][0]["page_limit"] = PAGE_LIMIT
    payload["task_name"] = f"html-to-pdf-pages-{idx}"
    return payload


@lru_cache(maxsize=None)
def _cached_html_payload() -> tuple[str, float]:
    body_parts = [f"<p>Page {p+1} - concurrency test</p><div style='page-break-after: always;'></div>" for p in range(DOC_PAGES)]
    html = "<html><body>" + "".join(body_parts) + "</body></html>"
    raw = html.encode("utf-8")
    b64 = base64.b64encode(raw).decode("ascii")
    size_mb = round(len(raw) / (1024 * 1024), 4)
    return b64, size_mb


@lru_cache(maxsize=None)
def _cached_docx_payload() -> tuple[str, float]:
    doc = Document()
    text = "汉" * DOC_CHARS_PER_PAGE
    for i in range(DOC_PAGES):
        p = doc.add_paragraph(f"DOCX混合测试 第{i+1}页 ")
        p.add_run(text)
        doc.add_page_break()
    with TemporaryDirectory() as tmpdir:
        path = os.path.join(tmpdir, "tmp.docx")
        doc.save(path)
        raw = Path(path).read_bytes()
    b64 = base64.b64encode(raw).decode("ascii")
    size_mb = round(len(raw) / (1024 * 1024), 4)
    return b64, size_mb


@lru_cache(maxsize=None)
def _cached_pptx_payload() -> tuple[str, float]:
    prs = Presentation()
    blank = prs.slide_layouts[6]
    text = "汉" * DOC_CHARS_PER_PAGE
    for i in range(DOC_PAGES):
        slide = prs.slides.add_slide(blank)
        tx_box = slide.shapes.add_textbox(50, 50, 600, 400)
        tf = tx_box.text_frame
        tf.text = f"PPTX混合测试 第{i+1}页 " + text
    with TemporaryDirectory() as tmpdir:
        path = os.path.join(tmpdir, "tmp.pptx")
        prs.save(path)
        raw = Path(path).read_bytes()
    b64 = base64.b64encode(raw).decode("ascii")
    size_mb = round(len(raw) / (1024 * 1024), 4)
    return b64, size_mb


@lru_cache(maxsize=None)
def _cached_wav_payload() -> tuple[str, float]:
    import io
    import struct
    import wave

    target_bytes = int(AUDIO_SIZE_MB * 1024 * 1024)
    buffer = io.BytesIO()
    with wave.open(buffer, "wb") as wav:
        wav.setnchannels(1)
        wav.setsampwidth(2)
        wav.setframerate(8000)
        # Each frame (sample) is 2 bytes for 16-bit mono. Compute frames to reach target size (approx).
        frames_needed = max(1, target_bytes // 2)
        chunk = struct.pack("<h", 0)
        wav.writeframes(chunk * frames_needed)
    raw = buffer.getvalue()[:target_bytes]
    b64 = base64.b64encode(raw).decode("ascii")
    size_mb = round(len(raw) / (1024 * 1024), 4)
    return b64, size_mb


def build_wav_payload(idx: int) -> Dict[str, Any]:
    b64, size_mb = _cached_wav_payload()
    return {
        "task_name": f"wav-to-mp3-{idx}",
        "priority": "normal",
        "files": [
            {
                "base64_data": b64,
                "source_format": "wav",
                "target_format": "mp3",
                "filename": f"silence-{idx}.wav",
                "size_mb": size_mb,
            }
        ],
    }


@lru_cache(maxsize=None)
def _cached_doc_payload() -> tuple[str, float]:
    # Generate doc via docx then convert to .doc using soffice once.
    b64_docx, _ = _cached_docx_payload()
    raw_docx = base64.b64decode(b64_docx.encode("ascii"))
    with TemporaryDirectory() as tmpdir:
        docx_path = Path(tmpdir) / "tmp.docx"
        doc_path = Path(tmpdir) / "tmp.doc"
        docx_path.write_bytes(raw_docx)
        import subprocess

        cmd = [
            "soffice",
            "--headless",
            "--convert-to",
            "doc",
            "--outdir",
            str(doc_path.parent),
            str(docx_path),
        ]
        subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        generated = doc_path.parent / f"{docx_path.stem}.doc"
        raw = generated.read_bytes()
    b64 = base64.b64encode(raw).decode("ascii")
    size_mb = round(len(raw) / (1024 * 1024), 4)
    return b64, size_mb


@lru_cache(maxsize=None)
def _cached_ppt_payload() -> tuple[str, float]:
    b64_pptx, _ = _cached_pptx_payload()
    raw_pptx = base64.b64decode(b64_pptx.encode("ascii"))
    with TemporaryDirectory() as tmpdir:
        pptx_path = Path(tmpdir) / "tmp.pptx"
        ppt_path = Path(tmpdir) / "tmp.ppt"
        pptx_path.write_bytes(raw_pptx)
        import subprocess

        cmd = [
            "soffice",
            "--headless",
            "--convert-to",
            "ppt",
            "--outdir",
            str(ppt_path.parent),
            str(pptx_path),
        ]
        subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        generated = ppt_path.parent / f"{pptx_path.stem}.ppt"
        raw = generated.read_bytes()
    b64 = base64.b64encode(raw).decode("ascii")
    size_mb = round(len(raw) / (1024 * 1024), 4)
    return b64, size_mb


@lru_cache(maxsize=None)
def _cached_gif_payload() -> tuple[str, float]:
    header = base64.b64decode("R0lGODlhAQABAPAAAAAAAAAAACH5BAEAAAAALAAAAAABAAEAAAICRAEAOw==")
    target_bytes = int(VIDEO_SIZE_MB * 1024 * 1024)
    raw = bytearray(header)
    if len(raw) < target_bytes:
        raw.extend(b"\0" * (target_bytes - len(raw)))
    raw = raw[:target_bytes]
    gif_b64 = base64.b64encode(raw).decode("ascii")
    size_mb = round(len(raw) / (1024 * 1024), 4)
    return gif_b64, size_mb


def build_gif_payload(idx: int) -> Dict[str, Any]:
    gif_b64, size_mb = _cached_gif_payload()
    return {
        "task_name": f"gif-to-mp4-{idx}",
        "priority": "normal",
        "files": [
            {
                "base64_data": gif_b64,
                "source_format": "gif",
                "target_format": "mp4",
                "filename": f"tiny-{idx}.gif",
                "size_mb": size_mb,
            }
        ],
    }


def build_doc_payload(idx: int) -> Dict[str, Any]:
    b64, size_mb = _cached_doc_payload()
    return {
        "task_name": f"doc-to-pdf-{idx}",
        "priority": "normal",
        "files": [
            {
                "base64_data": b64,
                "source_format": "doc",
                "target_format": "pdf",
                "filename": f"long-{idx}.doc",
                "size_mb": size_mb,
            }
        ],
    }


def build_docx_payload(idx: int) -> Dict[str, Any]:
    b64, size_mb = _cached_docx_payload()
    return {
        "task_name": f"docx-to-pdf-{idx}",
        "priority": "normal",
        "files": [
            {
                "base64_data": b64,
                "source_format": "docx",
                "target_format": "pdf",
                "filename": f"long-{idx}.docx",
                "size_mb": size_mb,
            }
        ],
    }


def build_ppt_payload(idx: int) -> Dict[str, Any]:
    b64, size_mb = _cached_ppt_payload()
    return {
        "task_name": f"ppt-to-pdf-{idx}",
        "priority": "normal",
        "files": [
            {
                "base64_data": b64,
                "source_format": "ppt",
                "target_format": "pdf",
                "filename": f"slides-{idx}.ppt",
                "size_mb": size_mb,
            }
        ],
    }


def build_pptx_payload(idx: int) -> Dict[str, Any]:
    b64, size_mb = _cached_pptx_payload()
    return {
        "task_name": f"pptx-to-pdf-{idx}",
        "priority": "normal",
        "files": [
            {
                "base64_data": b64,
                "source_format": "pptx",
                "target_format": "pdf",
                "filename": f"slides-{idx}.pptx",
                "size_mb": size_mb,
            }
        ],
    }


def build_doc_payload_pagelimit(idx: int) -> Dict[str, Any]:
    payload = build_doc_payload(idx)
    payload["files"][0]["page_limit"] = PAGE_LIMIT
    payload["task_name"] = f"doc-to-pdf-pages-{idx}"
    return payload


def build_docx_payload_pagelimit(idx: int) -> Dict[str, Any]:
    payload = build_docx_payload(idx)
    payload["files"][0]["page_limit"] = PAGE_LIMIT
    payload["task_name"] = f"docx-to-pdf-pages-{idx}"
    return payload


def build_ppt_payload_pagelimit(idx: int) -> Dict[str, Any]:
    payload = build_ppt_payload(idx)
    payload["files"][0]["page_limit"] = PAGE_LIMIT
    payload["task_name"] = f"ppt-to-pdf-pages-{idx}"
    return payload


def build_pptx_payload_pagelimit(idx: int) -> Dict[str, Any]:
    payload = build_pptx_payload(idx)
    payload["files"][0]["page_limit"] = PAGE_LIMIT
    payload["task_name"] = f"pptx-to-pdf-pages-{idx}"
    return payload


def build_requests() -> List[tuple[str, Dict[str, Any]]]:
    requests: List[tuple[str, Dict[str, Any]]] = []
    for i in range(REQUESTS_PER_SCENARIO):
        requests.append(("html(base64 50p)->pdf", build_html_payload(i)))
        requests.append(("html(base64 50p)->pdf (pages)", build_html_payload_pagelimit(i)))
        requests.append(("svg(base64)->png", build_svg_payload(i)))
        requests.append(("wav(base64)->mp3", build_wav_payload(i)))
        requests.append(("gif(base64)->mp4", build_gif_payload(i)))
        requests.append(("doc->pdf", build_doc_payload(i)))
        requests.append(("doc->pdf (pages)", build_doc_payload_pagelimit(i)))
        requests.append(("docx->pdf", build_docx_payload(i)))
        requests.append(("docx->pdf (pages)", build_docx_payload_pagelimit(i)))
        requests.append(("ppt->pdf", build_ppt_payload(i)))
        requests.append(("ppt->pdf (pages)", build_ppt_payload_pagelimit(i)))
        requests.append(("pptx->pdf", build_pptx_payload(i)))
        requests.append(("pptx->pdf (pages)", build_pptx_payload_pagelimit(i)))
    return requests


def summarize(results: List[dict[str, Any]]) -> None:
    total = len(results)
    ok = sum(1 for r in results if r["status"] == "ok")
    elapsed = [r["elapsed"] for r in results]
    p50 = sorted(elapsed)[int(0.5 * total)] if elapsed else 0
    p95 = sorted(elapsed)[int(0.95 * total) - 1] if elapsed else 0
    print(f"Summary: {ok}/{total} succeeded")
    print(f"Latency: p50={p50:.3f}s p95={p95:.3f}s max={max(elapsed) if elapsed else 0:.3f}s")
    by_scenario: dict[str, dict[str, Any]] = {}
    for r in results:
        rec = by_scenario.setdefault(r["scenario"], {"times": [], "ok": 0, "total": 0})
        rec["times"].append(r["elapsed"])
        rec["total"] += 1
        if r["status"] == "ok":
            rec["ok"] += 1

    sizes = describe_payloads()
    print("Per-scenario summary (with payload size/pages):")
    for scenario, rec in by_scenario.items():
        avg = sum(rec["times"]) / len(rec["times"]) if rec["times"] else 0
        size_key = SCENARIO_META.get(scenario)
        meta = sizes.get(size_key, {}) if size_key else {}
        size = meta.get("size_mb")
        pages = meta.get("pages")
        size_str = f"{size:.3f}MB" if isinstance(size, (int, float)) else "n/a"
        page_str = f", pages={pages}" if pages is not None else ""
        print(
            f"  {scenario}: {rec['total']} reqs, ok={rec['ok']}, avg={avg:.3f}s, size={size_str}{page_str}"
        )


def describe_payloads() -> dict[str, dict[str, float | int | None]]:
    # Trigger caches once and report sizes/pages.
    b64_html, size_html = _cached_html_payload()
    b64_svg, size_svg = _cached_svg_payload()
    b64_wav, size_wav = _cached_wav_payload()
    b64_gif, size_gif = _cached_gif_payload()
    b64_doc, size_doc = _cached_doc_payload()
    b64_docx, size_docx = _cached_docx_payload()
    b64_ppt, size_ppt = _cached_ppt_payload()
    b64_pptx, size_pptx = _cached_pptx_payload()

    return {
        "html->pdf": {"size_mb": size_html, "pages": DOC_PAGES},
        "svg->png": {"size_mb": size_svg, "pages": None},
        "wav->mp3": {"size_mb": size_wav, "pages": None},
        "gif->mp4": {"size_mb": size_gif, "pages": None},
        "doc->pdf": {"size_mb": size_doc, "pages": DOC_PAGES},
        "docx->pdf": {"size_mb": size_docx, "pages": DOC_PAGES},
        "ppt->pdf": {"size_mb": size_ppt, "pages": DOC_PAGES},
        "pptx->pdf": {"size_mb": size_pptx, "pages": DOC_PAGES},
    }


def main() -> int:
    requests = build_requests()
    results: List[dict[str, Any]] = []
    start = time.perf_counter()
    with ThreadPoolExecutor(max_workers=CONCURRENCY) as pool:
        future_to_name = {pool.submit(_post, payload, name): name for name, payload in requests}
        for fut in as_completed(future_to_name):
            res = fut.result()
            results.append(res)
            status = "OK" if res["status"] == "ok" else "FAIL"
            print(f"[{status}] {res['scenario']} ({res['elapsed']:.3f}s) -> {res['http']} {res['body'][:120]}")
    total_elapsed = time.perf_counter() - start
    print(f"Total wall time: {total_elapsed:.3f}s with concurrency={CONCURRENCY}")
    summarize(results)
    return 0 if all(r["status"] == "ok" for r in results) else 1


if __name__ == "__main__":
    raise SystemExit(main())
